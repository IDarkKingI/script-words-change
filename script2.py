import re
import os.path
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import datetime

def replace_words_in_docx(docx_path, dictionary, output_path, log_file):
    doc = Document(docx_path)
    replaced_words = set()
    for paragraph in doc.paragraphs:
        for old_word, new_word in dictionary.items():
            for run in paragraph.runs:
                replaced_text = re.sub(r'\b' + re.escape(old_word) + r'\b', new_word, run.text)
                if replaced_text != run.text:
                    replaced_words.add((old_word, new_word))
                run.text = replaced_text
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for old_word, new_word in dictionary.items():
                        pattern = r'\b' + re.escape(old_word) + r'\b'
                        replaced_text = re.sub(pattern, new_word, paragraph.text)
                        if replaced_text != paragraph.text:
                            replaced_words.add((old_word, new_word))
                        paragraph.text = replaced_text
    doc.save(output_path)
    if replaced_words:
        log_file.write(f"Replaced words in document: {os.path.basename(docx_path)}\n")
        log_file.write("Terms found for replacement:\n")
        for old_word, new_word in dictionary.items():
            log_file.write(f"{old_word} -> {new_word}\n")
        log_file.write("Terms where replacements were made:\n")
        for old_word, new_word in replaced_words:
            log_file.write(f"{old_word} -> {new_word}\n")
        log_file.write("\n")

def replace_words_in_pptx(pptx_path, dictionary, output_path, log_file):
    prs = Presentation(pptx_path)
    replaced_words = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for old_word, new_word in dictionary.items():
                    pattern = r'\b' + re.escape(old_word) + r'\b'
                    replaced_text = re.sub(pattern, new_word, shape.text)
                    if replaced_text != shape.text:
                        replaced_words.add((old_word, new_word))
                    shape.text = replaced_text
    prs.save(output_path)
    if replaced_words:
        log_file.write(f"Replaced words in presentation: {os.path.basename(pptx_path)}\n")
        log_file.write("Terms found for replacement:\n")
        for old_word, new_word in dictionary.items():
            log_file.write(f"{old_word} -> {new_word}\n")
        log_file.write("Terms where replacements were made:\n")
        for old_word, new_word in replaced_words:
            log_file.write(f"{old_word} -> {new_word}\n")
        log_file.write("\n")

def replace_words_in_xlsx(xlsx_path, dictionary, output_path, log_file):
    wb = load_workbook(xlsx_path)
    replaced_words = set()
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            for cell in row:
                for old_word, new_word in dictionary.items():
                    pattern = r'\b' + re.escape(old_word) + r'\b'
                    if cell.value and isinstance(cell.value, str):
                        replaced_text = re.sub(pattern, new_word, cell.value)
                        if replaced_text != cell.value:
                            replaced_words.add((old_word, new_word))
                        cell.value = replaced_text
    wb.save(output_path)
    if replaced_words:
        log_file.write(f"Replaced words in spreadsheet: {os.path.basename(xlsx_path)}\n")
        log_file.write("Terms found for replacement:\n")
        for old_word, new_word in dictionary.items():
            log_file.write(f"{old_word} -> {new_word}\n")
        log_file.write("Terms where replacements were made:\n")
        for old_word, new_word in replaced_words:
            log_file.write(f"{old_word} -> {new_word}\n")
        log_file.write("\n")

document_path = '444.docx'
presentation_path = '444.pptx'
spreadsheet_path = '444.xlsx'

replacements = {    #слова для замены
        "ЦК": "МК",
        "Центральная компания": "Маркетинговая компания",
        "ВЫСОЦКИЙ КОНСАЛТИНГ": "BUSINESS BOOSTER",
        "Visotsky Inc": "BUSINESS BOOSTER",
        "Visotsky Consulting": "BUSINESS BOOSTER",
        "VC Int": "BUSINESS BOOSTER",
        "ВК": "BUSINESS BOOSTER",
        "УК": "ОО",
        "Цель": "Видение ",
        "Замысел": "Миссия",
        "Инспекции": "Мониторинг",
        "Инспекции основ": "Мониторинг",
        "Tonnus": "BB Platform",
        "Статистика": "Метрика",
        "ИП": "Регламент",
        "ДУК": "Директива",
        "Проверочный список": "Чек-лист",
        "Оргсхема": "Оргструктура"
}

output_docx_path = '444_replaced.docx'
output_pptx_path = '444_replaced.pptx'
output_xlsx_path = '444_replaced.xlsx'

log_file_name = 'replacements_log.txt'

with open(log_file_name, 'w') as log_file:
    log_file.write(f"Replacements log - {datetime.datetime.now()}\n")
    if os.path.exists(document_path):
        replace_words_in_docx(document_path, replacements, output_docx_path, log_file)
    if os.path.exists(presentation_path):
        replace_words_in_pptx(presentation_path, replacements, output_pptx_path, log_file)
    if os.path.exists(spreadsheet_path):
        replace_words_in_xlsx(spreadsheet_path, replacements, output_xlsx_path, log_file)

print("Replacement process completed. Check the log file for details.")